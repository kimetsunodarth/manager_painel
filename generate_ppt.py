from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Colors
    color_bg = RGBColor(2, 6, 23)
    color_accent = RGBColor(56, 189, 248)
    color_text = RGBColor(248, 250, 252)
    color_secondary = RGBColor(148, 163, 184)
    color_purple = RGBColor(168, 85, 247)

    def apply_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_bg

    def add_title_slide(prs):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        apply_bg(slide)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "ANANIM MANAGER PAINEL"
        p.font.bold = True
        p.font.size = Pt(60)
        p.font.color.rgb = color_accent
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.add_paragraph()
        p2.text = "Enterprise Orchestration | v1.2.0"
        p2.font.size = Pt(24)
        p2.font.color.rgb = color_text
        p2.alignment = PP_ALIGN.CENTER

    def add_content_slide(prs, title_text, bullets, is_new=False):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        apply_bg(slide)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = color_purple if is_new else color_accent

        # Bullets
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for point in bullets:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = color_text
            p.space_after = Pt(15)

    # Slide 1
    add_title_slide(prs)

    # Slide 2: Problem
    add_content_slide(prs, "O Desafio", [
        "Gestão fragmentada entre Huawei Cloud e SAP Business One.",
        "Dificuldade em monitorar custos e ligar/desligar VMs em horários específicos.",
        "Falta de visibilidade centralizada de Backups (CBR) e Licenças.",
        "Operação manual demorada e propensa a erros humanos."
    ])

    # Slide 3: Solution
    add_content_slide(prs, "A Solução: Ananim Manager", [
        "Interface Unificada: Controle total de ECS, Backups e SAP Services.",
        "Segurança Enterprise: Isolamento granular de projetos por usuário.",
        "Automação Inteligente: Agendamentos precisos de infraestrutura.",
        "Auditoria Completa: Rastreabilidade de todas as ações de operadores."
    ])

    # Slide 4: v1.2.0 Features
    add_content_slide(prs, "Novidade v1.2.0: Monitor Externo", [
        "Compatibilidade com Cloud8 e outros sistemas externos.",
        "Monitoramento Passivo: Acompanha status real sem interferir nos comandos.",
        "Registro Automático de Horas Extras: Sistema detecta atividade imprevista.",
        "Inteligência de Sessão: Abertura e fechamento de sessões pelo monitorador.",
    ], is_new=True)

    # Slide 5: SAP Automation
    add_content_slide(prs, "SAP Control Center Mastery", [
        "Ativação de usuários de suporte totalmente automatizada.",
        "Retry Inteligente: Resiliência contra delays de carregamento do portal.",
        "Interface Simplificada: Operadores ativam suporte com um clique.",
        "Consolidação de dados de Backups e replicação CBR."
    ])

    # Slide 6: Strategic Value
    add_content_slide(prs, "Impacto Estratégico & ROI", [
        "Redução de 60% no tempo administrativo operacional.",
        "Payback Estimado: 2 meses (economia de horas de engenharia).",
        "Otimização de Custos Cloud: Garantia de desligamento fora do horário.",
        "Compliance: Auditoria em tempo real para padrões de segurança."
    ])

    # Slide 7: Conclusion
    add_content_slide(prs, "Pronto para Escalar", [
        "Dashboard Moderno e Responsivo.",
        "Arquitetura Resiliente integrada ao IIS.",
        "Ecossistema SAP & Huawei sob controle absoluto.",
        "Ananim Excellence v1.2.0 - Liderando a Automação."
    ])

    prs.save("APRESENTACAO.pptx")
    print("APRESENTACAO.pptx gerado com sucesso!")

if __name__ == "__main__":
    create_presentation()
